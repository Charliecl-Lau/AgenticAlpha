import pytest
from pathlib import Path
from src.renderer.slide_builder import build_pptx
from src.renderer.content_map import SlideSpec, SlideType


def _specs():
    return [
        SlideSpec(slide_type=SlideType.TITLE, title="Main Title", body="Subtitle text"),
        SlideSpec(slide_type=SlideType.FUNDAMENTALS, title="Fundamentals", body="Some body text"),
        SlideSpec(slide_type=SlideType.DISCLOSURE, title="Disclosure", body="AI used for X only."),
    ]


def test_build_pptx_creates_file(tmp_path):
    out = tmp_path / "test_deck.pptx"
    build_pptx(_specs(), str(out))
    assert out.exists()
    assert out.stat().st_size > 0


def test_build_pptx_has_correct_slide_count(tmp_path):
    from pptx import Presentation
    out = tmp_path / "test_deck.pptx"
    build_pptx(_specs(), str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == len(_specs())


def test_build_pptx_with_table_spec(tmp_path):
    from pptx import Presentation
    specs = [SlideSpec(
        slide_type=SlideType.FUNDAMENTALS,
        title="Margin Table",
        body="",
        table_rows=[["Metric", "CATL", "LGES"], ["Overseas GM", "31.4%", "N/A"]],
    )]
    out = tmp_path / "table_deck.pptx"
    build_pptx(specs, str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert any(shape.has_table for shape in prs.slides[0].shapes)
