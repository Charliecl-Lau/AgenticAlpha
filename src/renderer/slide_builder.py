import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.renderer.content_map import SlideSpec, SlideType

logger = logging.getLogger(__name__)

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)
_ACCENT = RGBColor(0x1F, 0x77, 0xB4)   # CATL blue
_DARK = RGBColor(0x22, 0x22, 0x22)


def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or _DARK
    return tx_box


def _render_title_slide(slide, spec: SlideSpec) -> None:
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 spec.title, font_size=36, bold=True, color=_ACCENT)
    _add_textbox(slide, Inches(1), Inches(4), Inches(11), Inches(0.8),
                 spec.body, font_size=20)


def _render_content_slide(slide, spec: SlideSpec) -> None:
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 spec.title, font_size=24, bold=True, color=_ACCENT)
    if spec.table_rows:
        rows = len(spec.table_rows)
        cols = len(spec.table_rows[0])
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2),
                                       Inches(12), Inches(0.4 * rows)).table
        for r_idx, row in enumerate(spec.table_rows):
            for c_idx, cell_text in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = cell_text
                cell.text_frame.paragraphs[0].runs[0].font.bold = (r_idx == 0)
    elif spec.chart_path:
        # Chart slides: embed image first, then caption below
        if Path(spec.chart_path).exists():
            slide.shapes.add_picture(spec.chart_path, Inches(0.5), Inches(1.2),
                                     Inches(12), Inches(4.8))
            if spec.body:
                _add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(1.0),
                             spec.body, font_size=12)
        else:
            logger.warning("Chart file not found, skipping embed: %s", spec.chart_path)
            if spec.body:
                _add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                             spec.body, font_size=16)
    elif spec.body:
        _add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                     spec.body, font_size=16)


def build_pptx(specs: list[SlideSpec], output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank_layout = prs.slide_layouts[6]   # blank layout
    for spec in specs:
        slide = prs.slides.add_slide(blank_layout)
        if spec.slide_type == SlideType.TITLE:
            _render_title_slide(slide, spec)
        else:
            _render_content_slide(slide, spec)
    prs.save(output_path)
